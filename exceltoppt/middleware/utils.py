import openai
import re
import os
import openpyxl
import yaml
import uuid

from openai import OpenAI

import pandas as pd
import streamlit as st
import numpy as np
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import aspose.slides as slides

from fastapi import UploadFile, File
from langchain.document_loaders import CSVLoader
from langchain.indexes import VectorstoreIndexCreator
from langchain.chains import RetrievalQA
from langchain.llms import OpenAI
from langchain.embeddings.openai import OpenAIEmbeddings

from functions import api_openai
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


#Load Environment Variables
load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")
os.environ["OPENAI_API_KEY"] = os.getenv("OPENAI_API_KEY")
current_woring_dir = os.path.dirname(os.path.realpath(__file__))


manifest_file_path = None
metrics = None
fiscal_year = None
fiscal_months = None
reporting_standards = None
csv_file_path = os.path.join(current_woring_dir, "data/Output", "output.csv")

cisco_logo = os.path.join("data/Templates", "cisco_logo.png")
cisco_slide_bg = os.path.join("data/Templates", "cisco_slide_bg.png")
cisco_toc_bg = os.path.join("data/Templates", "cisco_toc.png")
cisco_bg_metrics = os.path.join("data/Templates", "cisco_bg_metrics.png")


async def upload_file(file: UploadFile = File(...)):
    try:

        #Get File Type
        file_type = file.content_type

        #Get File Size
        file.file.seek(0,2)
        size = file.file.tell()
        file.file.seek(0)
        file_size = convert_bytes_to_human_readable(size)

        #Get File Name
        file_name = file.filename

        #Save Excel file to local drive
        os.makedirs("data/Excel", exist_ok=True)
        with open(f"data/Excel/{file.filename}", "wb") as f:
            f.write(file.file.read())

        #Excel/Workbook - get list of sheets
        sheets = []
        if file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or file_type == "application/vnd.ms-excel":
            file_path = f"data/Excel/{file.filename}"
            workbook = openpyxl.load_workbook(file_path, read_only=True)
            sheets = workbook.sheetnames
            
        return {"response": "Success", "file_size": file_size, "file_name": file_name, "file_type" : file_type, "sheet": sheets}
    except Exception as e:
        return {"response": "Error", "error_message": e}

def ask_questions(query, chain):
    try:
        response = chain({"question": query})
        return response['result']
    except Exception as e:
        return str(e)

def convert_bytes_to_human_readable(size_in_bytes):
    for unit in ['B', 'KB', 'MB', 'GB', 'TB']:
        if size_in_bytes < 1024.0:
            return f"{size_in_bytes:.2f} {unit}"
        size_in_bytes /= 1024.0

async def ask_your_doc(query: str, context: str):
    try:
        prompt = ""
        document_user_instructions = """
            Look at only my quesry and the context given below to provide the most relevant answer. If you are not sure or the relevant content is not available in the context 
            please mention that the document does not have the requested information to provide the appropriate answer.
            """
        prompt = document_user_instructions + " \n" + "query is given below \n " + query + "\n" + "context : " + context 
        response = api_openai(prompt)
        print(response)
        return str(response)
    except Exception as e:
        return str(e)
        
async def ask_openai(context: str, call_for: str):

    try:
        prompt = ""
        income_statements_prompt_template = """
            Identify all the dimensions like reporting standards, fiscal year, fiscal quarter, revenue, cost of sales, expenses, income products, service , business entity and all the metric amounts and and provide the result in the YAML file with the respective dimenstions and metrics in the following format given for one quarter and repeat for all quarters if the data is available in the context:
            Reporting Standards: 
             - GAAP
             - Non-GAAP
               - Fiscal Year: 
                - Fiscal Quarter:
                   - Quarter 1 (August)
                       - Revenue: 
                         - Product
                            - Metrics Amount
                         - Service
                            - Metrics Amount
                       - Cost of Sales
                         - Product
                            - Metrics Amount
                         - Service  
                            - Metrics Amount
                       - Gross Margins
                         - Metrics Amount
                       - Operating Income    
                         - Metrics Amount    
                       - Net Income    
                         - Metrics Amount 
                       - Net Income    
                         - Metrics Amount 
            """
        if(call_for == "generate_manifest"):
            income_statements_prompt_template = """
            Identify all the dimensions like reporting standards, fiscal year, fiscal quarter, revenue, cost of sales, expenses, income products, service , business entity and all the metric amounts and and provide the result in the YAML file with the respective dimenstions and metrics in the following format given for one quarter and repeat for all quarters if the data is available in the context:
            Reporting Standards: 
             - GAAP
             - Non-GAAP
               - Fiscal Year: 
                - Fiscal Quarter:
                   - Quarter 1 (August)
                       - Revenue: 
                         - Product
                            - Metrics Amount
                         - Service
                            - Metrics Amount
                       - Cost of Sales
                         - Product
                            - Metrics Amount
                         - Service  
                            - Metrics Amount
                       - Gross Margins
                         - Metrics Amount
                       - Operating Income    
                         - Metrics Amount    
                       - Net Income    
                         - Metrics Amount 
                       - Net Income    
                         - Metrics Amount 
            """
            prompt = income_statements_prompt_template + " from the context below" + f"\n\n{context}"
            #prompt = "Genenrate a manifest from the excel sheet in text format: " +f"\n\nExcel text: {context}"

        elif(call_for == "get_summary"):
            summary_prompt_template = """
            For the Excel sheet in text given below, which provides details about historic financial statements, Please provide a 5-point summary. In your response, give each point as a separate sentence, and label them as (1), (2), (3), (4), and (5). All the financial numbers must be in $ / dollars"""
            prompt = summary_prompt_template + f"\n\nExcel text: {context}"
         
        response = api_openai(prompt)
        if(call_for == "generate_manifest"):
            yaml_pattern = r'```yaml(.*?)```'
            #match = re.search(yaml_pattern, str(completion.choices[0].message.content), re.DOTALL)
            match = re.search(yaml_pattern, str(response), re.DOTALL)

            if match:
                yaml_content = match.group(1).strip()
                return yaml_content
            else:
               # return str(completion.choices[0].message.content)
               return str(response)
            
        if(call_for == "get_summary"):
            #return str(completion.choices[0].message.content)
            return str(response)
    except Exception as e:
        return str(e)

def excel_to_text(file_path: str, worksheet_name: str):

    try:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        worksheet_text_mapping = {}

        #get excel text from the provided worksheet
        if(worksheet_name is not None):
            worksheet = workbook[worksheet_name]
            text_content = ""
            for row in worksheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) if cell is not None else "" for cell in row)
                text_content += row_text + "\n"
            worksheet_text_mapping[worksheet_name] = text_content
        else:

        #get excel text from all worksheets
            for sheet_name in workbook.sheetnames:
            
                worksheet = workbook[sheet_name]
                text_content = ""

                for row in worksheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) if cell is not None else "" for cell in row)
                    text_content += row_text + "\n"

                worksheet_text_mapping[sheet_name] = text_content

        return worksheet_text_mapping
    except Exception as e:
        
        return str(e)

def listFilesInDirectory(directory_path: str):
    try:
        files_list = []
        for root, dirs, files in os.walk(directory_path):
            for file in files:
                file_name = os.path.basename(file)
                files_list.append(file_name)
        return files_list
    except Exception as e:
        return str(e)

async def generate_manifest(file_name: str, file_path: str, save_dir: str):

    directory_path = os.path.join(current_woring_dir, save_dir, file_name)
    if os.path.exists(directory_path):
        sheet_names = listFilesInDirectory(directory_path)
        return {"message": "Manifest generated successfully!!","fileName": file_name, "sheet_names": sheet_names, "status": 200}
    
    else:
    
        os.makedirs(os.path.join(current_woring_dir, save_dir, file_name), exist_ok=True)

        worksheets_text_mapping = excel_to_text(file_path, None)
        
        sheet_names = []
        for sheet_name, text_content in worksheets_text_mapping.items():
            
            
            response = await ask_openai(text_content, "generate_manifest")

            manifest_dir = os.path.join(current_woring_dir, save_dir, file_name, f"{sheet_name}.yaml")
            sheet_name = f"{sheet_name}.yaml"
            sheet_names.append(sheet_name)

            with open(manifest_dir, 'w') as manifest_file:
                manifest_file.write(response)
        
        return {"message": "Manifest generated successfully","fileName": file_name, "sheet_names": sheet_names, "status": 200}

def getcolumns(metric):
    return list(metric.values())[0]

def get_metric_row(df, keyword):
    row, column = None, None
    for i, col in enumerate(df.columns):
        for j, value in enumerate(df[col]):
            if str(value).strip().lower() == str(keyword).lower():
                row, column = j, i
                break
    return row

def get_common_row_and_column_2(df, keyword1, keyword2):
    common_row = None
    common_col = None

    for row in df.index:
        cleaned_row_values = [str(value).strip() for value in df.loc[row].values]
        if keyword1.strip() in cleaned_row_values and keyword2.strip() in cleaned_row_values:
            common_row = row
            break

    for col in df.columns:
        cleaned_col_values = [str(value).strip() for value in df[col].values]
        #print("cleaned_col_values: ", cleaned_col_values)
        if keyword1.strip() in cleaned_col_values and keyword2.strip() in cleaned_col_values:
            common_col = col
            break
    return common_row, common_col

def Excel_to_Dataframe_auto(uploaded_file, sheet_name, csv_file):
    # Load the Excel file
    wb = openpyxl.load_workbook(uploaded_file, data_only=True)
    ws = wb[sheet_name]

    # Convert the worksheet to a list of lists
    data = []
    for row in ws.iter_rows(values_only=True):
        data.append(row)

    # Remove empty rows
    data = [row for row in data if any(row)]

    # Create a DataFrame from the cleaned data
    header = data[0]  # Assuming the first row is the header row
    data = data[1:]  # Remove the header row
    df = pd.DataFrame(data, columns=header)

    # Save the DataFrame to a CSV file
    df.to_csv(csv_file, index=False)
    return df

def excel_to_csv(excel_file_path, selected_sheet):
    df = pd.read_excel(excel_file_path, sheet_name=selected_sheet, header=None)
    df = df.fillna(0)
    rid = uuid.uuid4()
    folderName = excel_file_path.split("/")[-1].split(".")[0]
    csv_file_name = selected_sheet + "_" + rid + ".csv" 
    os.makedirs(os.path.join("data/Manifest", folderName), exist_ok=True)
    rel_file_path = "data/Manifest" + folderName + csv_file_name
    df[selected_sheet].to_csv(rel_file_path)
    print("excel_to_csv " + rel_file_path)
    return csv_file_name

def Excel_to_dataframe(excel_file_path, manifest_file_path, selected_sheet):
    
    #load excel file into df
    df = pd.read_excel(excel_file_path, sheet_name=selected_sheet, header=None)
    df = df.fillna("")
    df = df.dropna(how='all')
    df = df.replace('\n', '', regex=True)

    #Extract YAML details
    with open(manifest_file_path, 'r') as yaml_file:
        data = yaml.safe_load(yaml_file)
    metrics = data['Metrics']
    fiscal_year = data['Fiscal Year']
    fiscal_months = data['Fiscal Months']
    reporting_standards = data['Reporting Standards']

    csv_df = pd.DataFrame()

    #constrcut csv from excel data
    for rs in reporting_standards:
        for month in fiscal_months:
            column = rs + " -[" + month.strip() + "]"
            column_values = []
            row_index = []
            row, col = get_common_row_and_column_2(df, rs, month)
            for m in metrics:
                if isinstance(m, str):
                    metric = m
                    m_row = get_metric_row(df, metric)
                    
                    if m_row is not None:
                        value = df.iat[m_row, col]
                        column_values.append(value)
                        row_index.append(metric.split('(')[0] if '(' in metric else metric)
                        #row_index.append(metric)
                if isinstance(m, dict):
                    for key, value in m.items():
                        metric = key
                        m_row = get_metric_row(df, metric)
                       
                        
                        if m_row is not None:
                            value = df.iat[m_row, col]
                            column_values.append(value)
                            row_index.append(metric.split('(')[0] if '(' in metric else metric)
                            #row_index.append(metric)
                        sub_metric = getcolumns(m)
                        for sm in sub_metric:
                            
                            sm_row = get_metric_row(df, sm)
                            
                            if sm_row is None:
                                val = df.iat[0, col]
                            else:
                                val = df.iat[sm_row, col]
                            column_values.append(val)
                            row_index.append((metric.split('(')[0] if '(' in metric else metric) + " - " + (sm.split('(')[0] if '(' in sm else sm))
                            #row_index.append(metric + " - " + sm)
            csv_df[column] = column_values
    csv_df.index = row_index
   # st.write(csv_df)
    return csv_df

def get_df_metric(df, rowindex, colindex):
    row = df[df.index.str.contains(rowindex)]
    #print("row: ", row)
    if not row.empty and colindex in df.columns:
        # Use .loc to access the specific cell at the intersection of row and col
        value = row.loc[row.index[0], colindex]
        return value

def formatNumber(number):
    
    result = float(number) * 100
    formatted_result = "{:.1f}%".format(result)
   
    return formatted_result

def add_footer(prs, slide, footer_text):

    slide_height = prs.slide_height
    footer = slide.shapes.add_textbox(Inches(0.25), slide_height -  Inches(1), Inches(5), Inches(0.5))
    text_frame_footer = footer.text_frame
    q = text_frame_footer.add_paragraph()
    q.text = footer_text
    q.font.size = Pt(12)
    q.font.color.rgb = RGBColor(128, 128, 128)
    q.alignment = PP_ALIGN.LEFT

def add_contents_to_column (prs, text_box, contents, start_number = 1):

    tf = text_box.text_frame
    tf.word_wrap = True

    for i, content in enumerate(contents, 1):
        p = tf.add_paragraph()
        p.text = f"{start_number}. {content}"
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(20) 
        run.font.color.rgb = RGBColor(255, 255, 255)  
        run.font.italic = True
        start_number += 1

def add_toc_slide(prs, contents):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide_width = prs.slide_width
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    pic = slide.shapes.add_picture(cisco_toc_bg, left, top, width, height)

    #logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(cisco_logo, imgleft, imgtop, width=image_width, height=image_height)

    title_text = "TABLE OF CONTENTS"
    left = Inches(0)
    top = Inches(0.5)  
    width = Inches(10)
    height = Inches(0.5) 
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    font = p.runs[0].font
    font.color.rgb = RGBColor(255, 255, 255)
    font.size = Pt(24)  
    font.bold = True
   

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5.5))

    if contents:
        total_items = len(contents)
        mid_point = total_items // 2

        add_contents_to_column(prs, left_box, contents[:mid_point])

        if total_items % 2 == 0:
            add_contents_to_column(prs, right_box, contents[mid_point:], start_number=mid_point + 1)
        else:
            add_contents_to_column(prs, right_box, contents[mid_point + 1:], start_number=mid_point + 2)   
    '''add_contents_to_column(prs, left_box, contents[:5])
    add_contents_to_column(prs, right_box, contents[5:], start_number=len(contents[:5]) + 1)'''
    add_footer(prs, slide, "© Cisco and/or its affiliates. All rights reserved. Cisco Confidential")

def add_summary(prs, summary):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide_width = prs.slide_width

    # Add a background picture
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    pic = slide.shapes.add_picture(cisco_slide_bg, left, top, width, height)

    #cisco logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(cisco_logo, imgleft, imgtop, width=image_width, height=image_height)

    # Title for the slide
    title_text = "SUMMARY"
    left = Inches(0)
    top = Inches(0.5)
    width = Inches(10)
    height = Inches(0.5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    font = p.runs[0].font
    font.color.rgb = RGBColor(255, 255, 255)
    font.size = Pt(24)
    font.bold = True


    left = Inches(0)
    top = Inches(1.25)  
    width = Inches(5)
    height = Inches(6)  
    text_box_left = slide.shapes.add_textbox(left, top, width, height)
    text_frame_left = text_box_left.text_frame
    summary_content = summary
    sentences = summary_content.split("\n")
    
    content = slide.shapes.add_textbox(Inches(1), Inches(2), slide_width - Inches(2), Inches(5))
    
    tf = content.text_frame  # Move this outside of the loop to create a single text frame for the entire content
    tf.word_wrap = True
    for point in sentences:
        p = tf.add_paragraph()
        p.text = point
        p.space_after = Inches(0.25)
        p.font.size = Pt(14)
        run = p.runs[0]
        run.font.color.rgb = RGBColor(255, 255, 255)
    
    add_footer(prs, slide, "© Cisco and/or its affiliates. All rights reserved. Cisco Confidential")

def add_title_slide_cisco_finance(prs, Title_Text, fiscal_year):
     #title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    background_color = RGBColor(13, 39, 77)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    prs.slides[0].background.fill.solid()
    prs.slides[0].background.fill.fore_color.rgb = background_color

    #cisco logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(cisco_logo, imgleft, imgtop, width=image_width, height=image_height)

    #header text
    left = width = height = Inches(1)
    top = Inches(2)
    wearecisco = slide.shapes.add_textbox(0, top, slide_width, height)
    wearecisco.text_frame.text = "#WEARECISCO"
    text_frame_wearecisco = wearecisco.text_frame
    paragraph = text_frame_wearecisco.paragraphs[0]
    font = paragraph.font
    font.name = 'Arial'
    font.size = Pt(35)  #
    font.color.rgb = RGBColor(0, 187, 235)
    font.underline = True
    paragraph.alignment = PP_ALIGN.LEFT

    #header text
    title = slide.shapes.add_textbox(0, Inches(3), slide_width, height)
    title.text_frame.text = Title_Text
    text_frame_title = title.text_frame
    paragraph = text_frame_title.paragraphs[0]
    font = paragraph.font
    font.name = 'Arial'
    font.size = Pt(30)
    font.color.rgb = RGBColor(0, 187, 235)
    paragraph.alignment = PP_ALIGN.LEFT
    sub_title = slide.shapes.add_textbox(0, Inches(4), slide_width, height)
    sub_title.text_frame.text = fiscal_year
    text_frame_sub_title = sub_title.text_frame
    paragraph = text_frame_sub_title.paragraphs[0]
    font = paragraph.font
    font.name = 'Arial'
    font.size = Pt(20)
    font.color.rgb = RGBColor(255, 255, 255)
    paragraph.alignment = PP_ALIGN.LEFT

    add_footer(prs, slide, "© Cisco and/or its affiliates. All rights reserved. Cisco Confidential")

def add_excel_contents (prs, metrics, rs, fiscal_months, sub_metrics, df):

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide_width = prs.slide_width
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    pic = slide.shapes.add_picture(cisco_bg_metrics, left, top, width, height)

    #logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(cisco_logo, imgleft, imgtop, width=image_width, height=image_height)

    title_text = metrics
    left = Inches(0)
    top = Inches(0.5)  
    width = Inches(10)
    height = Inches(0.5) 
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    font = p.runs[0].font
    font.name = 'Arial'
    font.color.rgb = RGBColor(0, 187, 235)
    font.size = Pt(24)  
    font.bold = True

    subtitle_text = rs
    left = Inches(0)
    top = Inches(0.95)  
    width = Inches(10)
    height = Inches(0.5) 
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = subtitle_text
    p.alignment = PP_ALIGN.CENTER
    font = p.runs[0].font
    font.name = 'Arial'
    font.color.rgb = RGBColor(255, 255, 255)
    font.size = Pt(20)  

    sub_text = "(in millions dollars, except per-share amounts)"
    left = Inches(0)
    top = Inches(1.5)  
    width = Inches(10)
    height = Inches(0.5) 
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = sub_text
    p.alignment = PP_ALIGN.CENTER
    font = p.runs[0].font
    font.name = 'Arial'
    font.color.rgb = RGBColor(255, 255, 255)
    font.size = Pt(12)  
   

    mleft = Inches(1)
    mtop = Inches(2)
    mwidth = Inches(4.5)
    mheight = Inches(5.5)

    i = 0
    for month in fiscal_months:
        if i < 2:
            mleft = Inches(0.5 + i * 5)
            mtop = mtop
            i += 1
        else:
            mleft = Inches(0.5 + (i - 2) * 5)
            mtop = Inches(5)
            i += 1

        txBox = slide.shapes.add_textbox(mleft, mtop, mwidth, mheight)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        q_text = f"Q{i} : "
        month_text = month

        # Add the Q text
        q_run = p.add_run()
        q_run.text = q_text
        q_run.font.bold = True
        q_run.font.color.rgb = RGBColor(0, 187, 235)

        # Add the month text
        month_run = p.add_run()
        month_run.text = month_text
        month_run.font.color.rgb = RGBColor(255, 255, 255)

        tf.word_wrap = True
         
        
        txBox = slide.shapes.add_textbox(mleft, mtop, mwidth, mheight)
        tf = txBox.text_frame
        p = tf.add_paragraph()

        # Create separate runs for different parts of the text
        q_text = f"Q{i} : "
        q_run.text = q_text
        q_run.font.bold = True
        month_text = month
        month_run.text = month_text

        # Add the entire text
        full_run = p.add_run()
        full_text = f"{q_text}{month_text}"
        full_run.text = full_text

        # Apply underline and blue color to the entire text
        full_run.font.underline = True
        full_run.font.color.rgb = RGBColor(0, 187, 235)

        

        tf.word_wrap = True
        for sm in sub_metrics:
            rowindex = metrics + " - " + (sm.split('(')[0] if '(' in sm else sm)
            colindex =  rs + " -[" + month.strip() + "]"
            val = get_df_metric(df, rowindex, colindex)
            
            try:
                val = formatNumber(val) if str(np.float64(val)).startswith("0.") else val
            except Exception as e:
                                val = ""
        
            sm_val_text = f"{sm.split('(')[0] if '(' in sm else sm} : {val}"
            p = tf.add_paragraph()
            p.text = sm_val_text
            p.font.color.rgb = RGBColor(255, 255, 255)

async def generate_summary(filePath: str, sheet_name: str):

    try:
       
        worksheet_text_mapping = excel_to_text(filePath, sheet_name)
        

        response = await ask_openai(worksheet_text_mapping[sheet_name], "get_summary")

        return {"summary": response, "status": 200}
    except Exception as e:
        return str(e)

def generate_cisco_presentation(excel_file_path, manifest_file_path, summary, selected_sheet):
    
    try:

        try:
            #Extract YAML details
            with open(manifest_file_path, 'r') as yaml_file:
                data = yaml.safe_load(yaml_file)
            metrics = data['Metrics']
            fiscal_year = data['Fiscal Year']
            fiscal_months = data['Fiscal Months']
            reporting_standards = data['Reporting Standards']
        except Exception as e:
            print ("Error while processing Manifest File: " + str(e))

        try:
            #Convert Excel to Datafeme/csv
            df = Excel_to_dataframe(excel_file_path, manifest_file_path, selected_sheet)
        except Exception as e:
            print("Error while converting Excel to Datafeme/csv: " + str(e))

        try:
            #create table of contents list
            contents = ["Summary"]

            for content in metrics:
                if isinstance(content, dict):
                    for key, value in content.items():
                        contents.append(key.split('(')[0] if '(' in key else key)


            contents.append("Miscellaneous")
        except Exception as e:
            print("Error while creating table of contents: " + str(e))
            contents = ["Summary", "Miscellaneous"]
    


        prs = Presentation()
        try:
            #title slide
            add_title_slide_cisco_finance(prs, "CISCO: Company Financial Overview", fiscal_year)
        except Exception as e:
            print("Error while creating title slide: " + str(e))
            add_title_slide_cisco_finance(prs, "CISCO: Company Financial Overview", "FY2023")

        try:
            #table of contents
            add_toc_slide(prs, contents)
        except Exception as e:
            print("Error while creating table of contents slide: " + str(e))

        try:
            #summary slide
            add_summary(prs, summary)
        except Exception as e:
            print("Error while creating summary slide: " + str(e) + "summary: " + str(summary))
            summary="""(1) line 1
            (2) line 2
            (3) line 3
            (4) line 4
            (5) line 5"""
            add_summary(prs, summary)

        try:
            #excel content slides
            for m in metrics:
                for rs in reporting_standards:
                    if isinstance(m, dict):
                        for key, value in m.items():
                            metric = key.split('(')[0] if '(' in key else key
                            sub_metric = getcolumns(m)
                            add_excel_contents(prs, metric, rs, fiscal_months, sub_metric, df)
        except Exception as e:
            print("Error while added excel contents: " + str(e))
                
        '''folderName = excel_file_path.split("/")[-1].split(".")[0]
        os.makedirs(os.path.join(current_woring_dir, "data/PowerPoints", folderName), exist_ok=True)
        #saving presentation
        ppt_name = str(selected_sheet).replace(" ", "") + "_presentation.pptx"'''
        ppt_path = os.path.join("data/PowerPoints", "cisco_presentation.pptx")
        pdf_path = os.path.join("data/pdfs", "cisco_presentation.pdf")
        prs.save(ppt_path)
        
        return ppt_path
    except Exception as e:
        return("Error while generating PowerPoints: " + str(e))

def generate_presentation(excel_file_path, manifest_file_path, summary, selected_sheet):
    
    try:
        #create table of contents list
        contents = ["Summary", "Company Overview", "Financial Overview", "Miscellaneous", "Conclusion"]

        prs = Presentation()
        try:
            #title slide
            add_title_slide_cisco_finance(prs, "CISCO: Company Financial Overview", "FY 2023")
        except Exception as e:
            print("Error while adding title slide: " + str(e))

        try:
            #table of contents
            add_toc_slide(prs, contents)
        except Exception as e:
            print("Error while adding table of content slide: " + str(e))

        try:
            #summary slide
            add_summary(prs, summary)
        except Exception as e:
            print("Error while adding summary slide: " + str(e))

        try:
            folderName = excel_file_path.split("/")[-1].split(".")[0]

            os.makedirs(os.path.join("data/PowerPoints", folderName), exist_ok=True)
            #saving presentation
            ppt_name = str(selected_sheet).replace(" ", "") + "_presentation.pptx"
            ppt_path = os.path.join("data/PowerPoints", folderName, ppt_name)
            prs.save(ppt_path)
            return ppt_path
        except Exception as e:
            folderName = "Excel_Sheets"
            os.makedirs(os.path.join("data/PowerPoints", folderName), exist_ok=True)
            #saving presentation
            ppt_name = str(selected_sheet).replace(" ", "") + "_presentation.pptx"
            pdf_name = str(selected_sheet).replace(" ", "") + "_presentation.pdf"
            ppt_path = os.path.join("data/PowerPoints", folderName, ppt_name)
            prs.save(ppt_path)
            prs.save(pdf_name)
            return ppt_path
    except Exception as e:
        return("Error while generating PowerPoints: " + str(e)) 

def ask_questions(query, chain):

    response = chain({"question": query})
    return response['result']

async def ask_csv_agent(excel_file_path, manifest_file_path, selected_sheet, question):

    try:
        df = Excel_to_dataframe(excel_file_path, manifest_file_path, selected_sheet)
        df.to_csv(csv_file_path, index = True)
        loader = CSVLoader(file_path=csv_file_path)
        document = loader.load()
        embeddings = OpenAIEmbeddings()
        index_creator = VectorstoreIndexCreator()
        docsearch = index_creator.from_loaders([loader])
        chain = RetrievalQA.from_chain_type(llm=OpenAI(), chain_type="stuff", retriever=docsearch.vectorstore.as_retriever(), input_key="question")
        response = ask_questions(question, chain)

        return {"response": response, "status": 200}
    except Exception as e:
        return str(e)








